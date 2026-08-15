import os
import io
import csv
from pathlib import Path
from typing import Generator
from dataclasses import dataclass

import fitz
from PIL import Image
import imagehash

try:
    import docx
    DOCX_OK = True
except ImportError:
    DOCX_OK = False

try:
    import openpyxl
    XLSX_OK = True
except ImportError:
    XLSX_OK = False

try:
    import pptx
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    from striprtf.striprtf import rtf_to_text
    RTF_OK = True
except ImportError:
    RTF_OK = False

BATCH_SIZE = 50
MIN_WORDS  = 2
IMAGE_DIR  = Path.home() / ".uniquest" / "extracted_images"


@dataclass
class TextChunk:
    content:     str
    page_number: int
    chunk_type:  str
    chunk_index: int


@dataclass
class ImageChunk:
    image_path:  str
    page_number: int
    image_index: int
    width:       int
    height:      int
    phash:       str
    ahash:       str
    dhash:       str


# ═════════════════════════════════════════════════════════════
def stream_pdf_text(file_path: str) -> Generator[list, None, None]:
    doc = fitz.open(file_path)
    chunk_index = 0
    batch = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text("text")
        for para in text.split("\n\n"):
            para = para.strip()
            if len(para.split()) >= MIN_WORDS:
                batch.append(TextChunk(para, page_num + 1, "paragraph", chunk_index))
                chunk_index += 1

        try:
            for table in page.find_tables():
                for row in table.extract():
                    row_texts = []
                    for cell in row:
                        if cell:
                            ct = str(cell).strip()
                            if len(ct.split()) >= MIN_WORDS:
                                batch.append(TextChunk(ct, page_num + 1, "cell", chunk_index))
                                chunk_index += 1
                            row_texts.append(ct)
                    row_combined = " | ".join(filter(None, row_texts))
                    if len(row_combined.split()) >= MIN_WORDS:
                        batch.append(TextChunk(row_combined, page_num + 1, "row", chunk_index))
                        chunk_index += 1
        except Exception:
            pass

        if (page_num + 1) % BATCH_SIZE == 0:
            yield batch
            batch = []

    if batch:
        yield batch
    doc.close()


def stream_pdf_images(file_path: str, project_id: int) -> Generator[list, None, None]:
    out_dir = IMAGE_DIR / str(project_id)
    out_dir.mkdir(parents=True, exist_ok=True)

    doc = fitz.open(file_path)
    image_index = 0
    batch = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        for img in page.get_images(full=True):
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
                img_bytes  = base_image["image"]
                ext        = base_image["ext"]
                pil_img    = Image.open(io.BytesIO(img_bytes)).convert("RGB")

                if pil_img.width < 50 or pil_img.height < 50:
                    continue

                img_path = out_dir / f"img_{page_num}_{image_index}.{ext}"
                pil_img.save(str(img_path))

                batch.append(ImageChunk(
                    image_path=str(img_path),
                    page_number=page_num + 1,
                    image_index=image_index,
                    width=pil_img.width,
                    height=pil_img.height,
                    phash=str(imagehash.phash(pil_img)),
                    ahash=str(imagehash.average_hash(pil_img)),
                    dhash=str(imagehash.dhash(pil_img)),
                ))
                image_index += 1
            except Exception:
                continue

        if (page_num + 1) % BATCH_SIZE == 0:
            yield batch
            batch = []

    if batch:
        yield batch
    doc.close()


# ═════════════════════════════════════════════════════════════
def extract_docx_text(file_path: str) -> list:
    if not DOCX_OK:
        return []
    chunks = []
    idx = 0
    try:
        document = docx.Document(file_path)
        for para in document.paragraphs:
            text = para.text.strip()
            if len(text.split()) >= MIN_WORDS:
                chunks.append(TextChunk(text, 1, "paragraph", idx))
                idx += 1
        for table in document.tables:
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    ct = cell.text.strip()
                    if len(ct.split()) >= MIN_WORDS:
                        chunks.append(TextChunk(ct, 1, "cell", idx))
                        idx += 1
                    row_texts.append(ct)
                row_combined = " | ".join(filter(None, row_texts))
                if len(row_combined.split()) >= MIN_WORDS:
                    chunks.append(TextChunk(row_combined, 1, "row", idx))
                    idx += 1
    except Exception:
        pass
    return chunks


def extract_xlsx_text(file_path: str) -> list:
    if not XLSX_OK:
        return []
    chunks = []
    idx = 0
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_texts = []
                for cell in row:
                    if cell is not None:
                        ct = str(cell).strip()
                        if len(ct.split()) >= MIN_WORDS:
                            chunks.append(TextChunk(ct, 1, "cell", idx))
                            idx += 1
                        row_texts.append(ct)
                row_combined = " | ".join(filter(None, row_texts))
                if len(row_combined.split()) >= MIN_WORDS:
                    chunks.append(TextChunk(row_combined, 1, "row", idx))
                    idx += 1
        wb.close()
    except Exception:
        pass
    return chunks


def extract_csv_text(file_path: str) -> list:
    chunks = []
    idx = 0
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                row_texts = []
                for cell in row:
                    ct = cell.strip()
                    if len(ct.split()) >= MIN_WORDS:
                        chunks.append(TextChunk(ct, 1, "cell", idx))
                        idx += 1
                    row_texts.append(ct)
                row_combined = " | ".join(filter(None, row_texts))
                if len(row_combined.split()) >= MIN_WORDS:
                    chunks.append(TextChunk(row_combined, 1, "row", idx))
                    idx += 1
    except Exception:
        pass
    return chunks


def extract_pptx_text(file_path: str) -> list:
    if not PPTX_OK:
        return []
    chunks = []
    idx = 0
    try:
        prs = pptx.Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if len(text.split()) >= MIN_WORDS:
                            chunks.append(TextChunk(text, slide_num, "paragraph", idx))
                            idx += 1
                        slide_text.append(text)
            combined = " ".join(filter(None, slide_text))
            if len(combined.split()) >= MIN_WORDS:
                chunks.append(TextChunk(combined, slide_num, "slide", idx))
                idx += 1
    except Exception:
        pass
    return chunks


def extract_txt_text(file_path: str) -> list:
    chunks = []
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        idx = 0
        for para in content.split("\n\n"):
            para = para.strip()
            if len(para.split()) >= MIN_WORDS:
                chunks.append(TextChunk(para, 1, "paragraph", idx))
                idx += 1
    except Exception:
        pass
    return chunks


def extract_rtf_text(file_path: str) -> list:
    if not RTF_OK:
        return []
    chunks = []
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            raw = f.read()
        plain = rtf_to_text(raw)
        idx = 0
        for para in plain.split("\n\n"):
            para = para.strip()
            if len(para.split()) >= MIN_WORDS:
                chunks.append(TextChunk(para, 1, "paragraph", idx))
                idx += 1
    except Exception:
        pass
    return chunks


# ═════════════════════════════════════════════════════════════
def get_file_type(file_path: str) -> str:
    return Path(file_path).suffix.lower().lstrip(".")


def extract_text_streaming(file_path: str):
    ft = get_file_type(file_path)
    if ft == "pdf":
        yield from stream_pdf_text(file_path)
    else:
        if ft == "docx":
            chunks = extract_docx_text(file_path)
        elif ft == "xlsx":
            chunks = extract_xlsx_text(file_path)
        elif ft == "csv":
            chunks = extract_csv_text(file_path)
        elif ft == "pptx":
            chunks = extract_pptx_text(file_path)
        elif ft in ("txt", "text"):
            chunks = extract_txt_text(file_path)
        elif ft == "rtf":
            chunks = extract_rtf_text(file_path)
        else:
            chunks = []
        if chunks:
            yield chunks


def extract_images_streaming(file_path: str, project_id: int):
    ft = get_file_type(file_path)
    if ft == "pdf":
        yield from stream_pdf_images(file_path, project_id)
try:
    import pytesseract
    from PIL import Image as PILImage
    # Auto-detect Tesseract location
    import os as _os
    for _p in (
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    ):
        if _os.path.exists(_p):
            pytesseract.pytesseract.tesseract_cmd = _p
            break
    OCR_AVAILABLE = True
except Exception:
    OCR_AVAILABLE = False