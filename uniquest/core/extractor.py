import os
import re
import csv
import shutil
from pathlib import Path
from typing import List, Tuple, Optional

from database.db import get_connection, get_db_path
from database.models import (
    TextChunk, ExtractedImage,
    IMAGE_EXTENSIONS, DOCUMENT_EXTENSIONS
)


# ─────────────────────────────────────────────
#  IMAGES STORAGE FOLDER
# ─────────────────────────────────────────────
def get_images_dir(project_id: int) -> Path:
    base = Path(get_db_path()).parent / "extracted_images" / str(project_id)
    base.mkdir(parents=True, exist_ok=True)
    return base


# ─────────────────────────────────────────────
#  TEXT CLEANING
# ─────────────────────────────────────────────
def clean_text(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r'\r\n|\r', '\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = text.strip()
    return text


def make_chunks(text: str, min_words: int = 2) -> List[Tuple[str, str]]:
    """
    Multi-level chunking:
      - paragraphs (>= 8 words)
      - lines      (>= 3 words)
      - short entries (>= 2 words, like table cells)
    Returns list of (chunk_text, chunk_type)
    """
    if not text:
        return []

    chunks: List[Tuple[str, str]] = []
    seen: set = set()

    def add(txt: str, ctype: str):
        txt = txt.strip()
        if not txt:
            return
        wc = len(txt.split())
        if wc < min_words:
            return
        key = txt.lower()
        if key in seen:
            return
        seen.add(key)
        chunks.append((txt, ctype))

    # ── 1. Paragraph-level chunks (>=8 words) ──
    for para in text.split('\n\n'):
        para = para.strip()
        if para and len(para.split()) >= 8:
            add(para, "paragraph")

    # ── 2. Line-level chunks (>=3 words) ──
    for line in text.split('\n'):
        line = line.strip()
        if line and len(line.split()) >= 3:
            add(line, "line")

    # ── 3. Short entries (2 words) — for table cells, lists ──
    for line in text.split('\n'):
        line = line.strip()
        parts = re.split(r'\s*[|,;\t]\s*', line)
        for part in parts:
            part = part.strip()
            if not part:
                continue
            wc = len(part.split())
            if wc >= min_words and wc < 8:
                # Only keep if it looks meaningful
                if len(part) >= 3 and not part.isdigit():
                    add(part, "cell")

    return chunks


def chunks_to_objects(
    chunks_data: List[Tuple[str, str]],
    file_id: int,
    project_id: int,
    page_number: int = 0,
) -> List[TextChunk]:
    """Convert list of (text, type) into TextChunk objects"""
    objects = []
    for idx, (content, ctype) in enumerate(chunks_data):
        objects.append(TextChunk(
            file_id=file_id,
            project_id=project_id,
            chunk_index=idx,
            content=content,
            page_number=page_number,
            chunk_type=ctype,
            word_count=len(content.split()),
        ))
    return objects


# ─────────────────────────────────────────────
#  PDF EXTRACTOR
# ─────────────────────────────────────────────
def extract_from_pdf(file_id, project_id, file_path):
    text_chunks = []
    image_paths = []

    try:
        import fitz
    except ImportError:
        print("PyMuPDF not installed.")
        return text_chunks, image_paths

    try:
        doc = fitz.open(file_path)
        images_dir = get_images_dir(project_id)
        global_idx = 0

        for page_num, page in enumerate(doc, start=1):
            raw_text = page.get_text("text")
            cleaned = clean_text(raw_text)
            chunks_data = make_chunks(cleaned, min_words=2)

            for content, ctype in chunks_data:
                text_chunks.append(TextChunk(
                    file_id=file_id,
                    project_id=project_id,
                    chunk_index=global_idx,
                    content=content,
                    page_number=page_num,
                    chunk_type=ctype,
                    word_count=len(content.split()),
                ))
                global_idx += 1

            # Also extract tables specifically
            try:
                tables = page.find_tables()
                for tbl_idx, table in enumerate(tables):
                    rows = table.extract()
                    for row_idx, row in enumerate(rows):
                        for cell in row:
                            if cell is None:
                                continue
                            cell_txt = str(cell).strip()
                            if not cell_txt:
                                continue
                            wc = len(cell_txt.split())
                            if wc >= 1 and len(cell_txt) >= 3:
                                text_chunks.append(TextChunk(
                                    file_id=file_id,
                                    project_id=project_id,
                                    chunk_index=global_idx,
                                    content=cell_txt,
                                    page_number=page_num,
                                    chunk_type="table_cell",
                                    word_count=wc,
                                ))
                                global_idx += 1
            except Exception:
                pass  # Tables not available in older PyMuPDF

            # Images
            image_list = page.get_images(full=True)
            for img_idx, img_info in enumerate(image_list):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    ext = base_image["ext"]
                    img_name = f"f{file_id}_p{page_num}_i{img_idx}.{ext}"
                    img_path = images_dir / img_name
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)
                    image_paths.append(str(img_path))
                except Exception as e:
                    print(f"  Image extract error: {e}")

        doc.close()
    except Exception as e:
        print(f"PDF extraction error [{file_path}]: {e}")

    return text_chunks, image_paths


# ─────────────────────────────────────────────
#  DOCX EXTRACTOR
# ─────────────────────────────────────────────
def extract_from_docx(file_id, project_id, file_path):
    text_chunks = []
    image_paths = []

    try:
        from docx import Document
    except ImportError:
        print("python-docx not installed.")
        return text_chunks, image_paths

    try:
        doc = Document(file_path)
        images_dir = get_images_dir(project_id)
        global_idx = 0

        # Paragraph text
        full_text = "\n\n".join(
            p.text for p in doc.paragraphs if p.text.strip()
        )
        cleaned = clean_text(full_text)
        chunks_data = make_chunks(cleaned, min_words=2)
        for content, ctype in chunks_data:
            text_chunks.append(TextChunk(
                file_id=file_id,
                project_id=project_id,
                chunk_index=global_idx,
                content=content,
                page_number=0,
                chunk_type=ctype,
                word_count=len(content.split()),
            ))
            global_idx += 1

        # Table cells
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    cell_txt = cell.text.strip()
                    if not cell_txt or len(cell_txt) < 3:
                        continue
                    text_chunks.append(TextChunk(
                        file_id=file_id,
                        project_id=project_id,
                        chunk_index=global_idx,
                        content=cell_txt,
                        page_number=0,
                        chunk_type="table_cell",
                        word_count=len(cell_txt.split()),
                    ))
                    global_idx += 1

        # Embedded images
        import zipfile
        with zipfile.ZipFile(file_path, 'r') as z:
            media_files = [
                f for f in z.namelist()
                if f.startswith("word/media/")
            ]
            for mf in media_files:
                img_name = f"f{file_id}_{Path(mf).name}"
                img_path = images_dir / img_name
                with z.open(mf) as src, open(img_path, "wb") as dst:
                    dst.write(src.read())
                image_paths.append(str(img_path))

    except Exception as e:
        print(f"DOCX extraction error [{file_path}]: {e}")

    return text_chunks, image_paths


# ─────────────────────────────────────────────
#  TXT / RTF EXTRACTOR
# ─────────────────────────────────────────────
def extract_from_txt(file_id, project_id, file_path):
    text_chunks = []

    try:
        ext = file_path.rsplit(".", 1)[-1].lower()

        if ext == "rtf":
            try:
                from striprtf.striprtf import rtf_to_text
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    raw = rtf_to_text(f.read())
            except ImportError:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    raw = f.read()
        else:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                raw = f.read()

        cleaned = clean_text(raw)
        chunks_data = make_chunks(cleaned, min_words=2)
        for idx, (content, ctype) in enumerate(chunks_data):
            text_chunks.append(TextChunk(
                file_id=file_id,
                project_id=project_id,
                chunk_index=idx,
                content=content,
                page_number=0,
                chunk_type=ctype,
                word_count=len(content.split()),
            ))

    except Exception as e:
        print(f"TXT/RTF extraction error [{file_path}]: {e}")

    return text_chunks, []


# ─────────────────────────────────────────────
#  XLSX / XLS / CSV EXTRACTOR
# ─────────────────────────────────────────────
def extract_from_spreadsheet(file_id, project_id, file_path):
    text_chunks = []
    ext = file_path.rsplit(".", 1)[-1].lower()
    global_idx = 0

    try:
        if ext == "csv":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                reader = csv.reader(f)
                for row_idx, row in enumerate(reader):
                    # Each cell as a chunk
                    for cell in row:
                        cell = cell.strip()
                        if cell and len(cell) >= 3 and not cell.isdigit():
                            text_chunks.append(TextChunk(
                                file_id=file_id,
                                project_id=project_id,
                                chunk_index=global_idx,
                                content=cell,
                                page_number=0,
                                chunk_type="table_cell",
                                word_count=len(cell.split()),
                            ))
                            global_idx += 1
                    # Full row as combined chunk
                    line = " | ".join(
                        c.strip() for c in row if c.strip()
                    )
                    if line and len(line.split()) >= 2:
                        text_chunks.append(TextChunk(
                            file_id=file_id,
                            project_id=project_id,
                            chunk_index=global_idx,
                            content=line,
                            page_number=0,
                            chunk_type="row",
                            word_count=len(line.split()),
                        ))
                        global_idx += 1
        else:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    cells_str = []
                    for cell in row:
                        if cell is None:
                            continue
                        cell_txt = str(cell).strip()
                        if not cell_txt:
                            continue
                        cells_str.append(cell_txt)
                        # Each cell as a chunk
                        if len(cell_txt) >= 3 and not cell_txt.isdigit():
                            text_chunks.append(TextChunk(
                                file_id=file_id,
                                project_id=project_id,
                                chunk_index=global_idx,
                                content=cell_txt,
                                page_number=0,
                                chunk_type="table_cell",
                                word_count=len(cell_txt.split()),
                            ))
                            global_idx += 1
                    # Row as combined chunk
                    if cells_str:
                        line = " | ".join(cells_str)
                        if len(line.split()) >= 2:
                            text_chunks.append(TextChunk(
                                file_id=file_id,
                                project_id=project_id,
                                chunk_index=global_idx,
                                content=line,
                                page_number=0,
                                chunk_type="row",
                                word_count=len(line.split()),
                            ))
                            global_idx += 1
            wb.close()

    except Exception as e:
        print(f"Spreadsheet extraction error [{file_path}]: {e}")

    return text_chunks, []


# ─────────────────────────────────────────────
#  PPTX EXTRACTOR
# ─────────────────────────────────────────────
def extract_from_pptx(file_id, project_id, file_path):
    text_chunks = []
    image_paths = []

    try:
        from pptx import Presentation
    except ImportError:
        print("python-pptx not installed.")
        return text_chunks, image_paths

    try:
        prs = Presentation(file_path)
        images_dir = get_images_dir(project_id)
        global_idx = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            slide_text.append(line)
                            # Each line as chunk if >=2 words
                            if len(line.split()) >= 2 and len(line) >= 3:
                                text_chunks.append(TextChunk(
                                    file_id=file_id,
                                    project_id=project_id,
                                    chunk_index=global_idx,
                                    content=line,
                                    page_number=slide_num,
                                    chunk_type="line",
                                    word_count=len(line.split()),
                                ))
                                global_idx += 1

                if shape.shape_type == 13:
                    try:
                        img_blob = shape.image.blob
                        ext = shape.image.ext
                        img_name = f"f{file_id}_s{slide_num}_sh{shape.shape_id}.{ext}"
                        img_path = images_dir / img_name
                        with open(img_path, "wb") as f:
                            f.write(img_blob)
                        image_paths.append(str(img_path))
                    except Exception as e:
                        print(f"  PPTX image error: {e}")

            # Full slide as combined chunk
            if slide_text:
                combined = "\n".join(slide_text)
                cleaned = clean_text(combined)
                if len(cleaned.split()) >= 8:
                    text_chunks.append(TextChunk(
                        file_id=file_id,
                        project_id=project_id,
                        chunk_index=global_idx,
                        content=cleaned,
                        page_number=slide_num,
                        chunk_type="slide",
                        word_count=len(cleaned.split()),
                    ))
                    global_idx += 1

    except Exception as e:
        print(f"PPTX extraction error [{file_path}]: {e}")

    return text_chunks, image_paths


# ─────────────────────────────────────────────
#  IMAGE FILE EXTRACTOR
# ─────────────────────────────────────────────
def extract_from_image(file_id, project_id, file_path):
    image_paths = []

    try:
        images_dir = get_images_dir(project_id)
        ext = file_path.rsplit(".", 1)[-1].lower()
        img_name = f"f{file_id}_standalone.{ext}"
        dest = images_dir / img_name

        if ext == "svg":
            return [], []

        shutil.copy2(file_path, dest)
        image_paths.append(str(dest))

    except Exception as e:
        print(f"Image copy error [{file_path}]: {e}")

    return [], image_paths


# ─────────────────────────────────────────────
#  MASTER EXTRACTOR
# ─────────────────────────────────────────────
def extract_file(file_id, project_id, file_path, file_type):
    ext = file_type.lower()

    if ext == "pdf":
        return extract_from_pdf(file_id, project_id, file_path)
    elif ext in ("docx", "doc"):
        return extract_from_docx(file_id, project_id, file_path)
    elif ext in ("txt", "rtf"):
        return extract_from_txt(file_id, project_id, file_path)
    elif ext in ("xlsx", "xls", "csv"):
        return extract_from_spreadsheet(file_id, project_id, file_path)
    elif ext in ("pptx", "ppt"):
        return extract_from_pptx(file_id, project_id, file_path)
    elif ext in ("jpg", "jpeg", "png", "bmp", "tiff", "tif", "webp", "gif"):
        return extract_from_image(file_id, project_id, file_path)
    else:
        print(f"Unsupported file type: {ext}")
        return [], []


# ─────────────────────────────────────────────
#  SAVE TO DATABASE
# ─────────────────────────────────────────────
def save_text_chunks(chunks):
    if not chunks:
        return 0
    conn = get_connection()
    cursor = conn.cursor()
    saved = 0
    for chunk in chunks:
        try:
            cursor.execute("""
                INSERT INTO text_chunks
                    (file_id, project_id, chunk_index,
                     content, page_number, chunk_type, word_count)
                VALUES (?, ?, ?, ?, ?, ?, ?)
            """, (
                chunk.file_id,
                chunk.project_id,
                chunk.chunk_index,
                chunk.content,
                chunk.page_number,
                chunk.chunk_type,
                chunk.word_count,
            ))
            saved += 1
        except Exception as e:
            print(f"  Chunk save error: {e}")
    conn.commit()
    conn.close()
    return saved


def save_extracted_images(file_id, project_id, image_paths):
    if not image_paths:
        return 0

    try:
        from PIL import Image
        import imagehash
    except ImportError:
        print("Pillow / imagehash not installed.")
        return 0

    conn = get_connection()
    cursor = conn.cursor()
    saved = 0

    for idx, img_path in enumerate(image_paths):
        try:
            img = Image.open(img_path).convert("RGB")
            ph = str(imagehash.phash(img))
            ah = str(imagehash.average_hash(img))
            dh = str(imagehash.dhash(img))
            w, h = img.size
            fsize = os.path.getsize(img_path)

            cursor.execute("""
                INSERT INTO extracted_images
                    (file_id, project_id, image_index,
                     stored_path, page_number,
                     width, height, phash, ahash, dhash, file_size)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """, (
                file_id, project_id, idx,
                img_path, 0,
                w, h, ph, ah, dh, fsize,
            ))
            saved += 1
        except Exception as e:
            print(f"  Image hash error [{img_path}]: {e}")

    conn.commit()
    conn.close()
    return saved


def update_file_status(file_id, status, text_count=0, image_count=0, error=None):
    conn = get_connection()
    cursor = conn.cursor()
    cursor.execute("""
        UPDATE files SET
            status = ?,
            text_extracted = ?,
            images_extracted = ?,
            error_message = ?,
            processed_at = CURRENT_TIMESTAMP
        WHERE id = ?
    """, (status, text_count, image_count, error, file_id))
    conn.commit()
    conn.close()


# ─────────────────────────────────────────────
#  FULL EXTRACTION PIPELINE FOR ONE FILE
# ─────────────────────────────────────────────
def process_file_extraction(file_id, project_id, file_path, file_type):
    print(f"  Extracting: {Path(file_path).name}")

    try:
        text_chunks, image_paths = extract_file(
            file_id, project_id, file_path, file_type
        )
        text_saved = save_text_chunks(text_chunks)
        image_saved = save_extracted_images(file_id, project_id, image_paths)
        update_file_status(file_id, "done", text_saved, image_saved)
        print(f"    OK: {text_saved} chunks, {image_saved} images")
        return text_saved, image_saved

    except Exception as e:
        error_msg = str(e)
        print(f"    ERROR: {error_msg}")
        update_file_status(file_id, "error", error=error_msg)
        return 0, 0